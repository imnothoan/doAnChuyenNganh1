from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

try:
    from scripts._bootstrap import ensure_repo_root_on_path
except ModuleNotFoundError:
    from _bootstrap import ensure_repo_root_on_path

ensure_repo_root_on_path(__file__)

from src.utils.config import CFG


ROOT = Path(__file__).resolve().parents[1]
SLIDES_DIR = ROOT / "slides"
OUT = SLIDES_DIR / "major_project_1_defense.pptx"

BLACK = RGBColor(20, 20, 20)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(92, 92, 92)
LIGHT_GRAY = RGBColor(238, 238, 238)
MID_GRAY = RGBColor(170, 170, 170)
DARK_GRAY = RGBColor(55, 55, 55)


def set_run(run, size=22, bold=False, color=BLACK):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def textbox(slide, text, x, y, w, h, size=22, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def bullet_list(slide, items, x, y, w, h, size=20, color=BLACK, gap=4):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(gap)
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_header(slide, title, number):
    textbox(slide, title, 0.55, 0.26, 10.7, 0.42, size=19, bold=True, color=BLACK)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.78), Inches(12.25), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BLACK
    line.line.fill.background()
    textbox(slide, f"{number:02d}", 12.15, 0.24, 0.65, 0.3, size=11, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)


def add_footer(slide):
    textbox(
        slide,
        "Machine Learning-based News Reliability Assessment | Major Project 1",
        0.55,
        7.12,
        11.6,
        0.22,
        size=9,
        color=GRAY,
    )


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    textbox(slide, "Machine Learning-based", 0.72, 1.18, 11.9, 0.55, size=28, bold=True)
    textbox(slide, "Visual Tool for News Reliability Assessment", 0.72, 1.74, 11.9, 0.82, size=34, bold=True)
    textbox(slide, "A Vietnamese NLP/ML web application for reliability screening, explanation, and feedback collection", 0.76, 2.82, 11.2, 0.52, size=18, color=DARK_GRAY)

    labels = [
        ("Streamlit UI", 0.78),
        ("NLP/ML Core", 3.0),
        ("Supabase/PostgreSQL", 5.22),
        ("Explainable Output", 7.85),
    ]
    for label, x in labels:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.18), Inches(1.95), Inches(0.62))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.color.rgb = BLACK
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_run(r, size=14, bold=True)
    textbox(slide, "Defense deck | 04/06/2026", 0.78, 6.55, 5, 0.28, size=13, color=GRAY)


def add_objectives_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Research Objectives and Contributions", number)
    textbox(slide, "Objectives", 0.72, 1.18, 5.3, 0.38, size=21, bold=True)
    bullet_list(
        slide,
        [
            "Build a web-based tool for Vietnamese news reliability screening.",
            "Train and compare supervised ML baselines for fake/clickbait detection.",
            "Visualize risk, model signals, suspicious terms, and prediction history.",
            "Collect user feedback as a future retraining source.",
        ],
        0.9,
        1.75,
        5.55,
        3.6,
        size=17,
    )
    textbox(slide, "Main contributions", 6.72, 1.18, 5.3, 0.38, size=21, bold=True)
    bullet_list(
        slide,
        [
            "Reproducible data pipeline: download, clean, split, train, evaluate.",
            "Four-model benchmark with a metric-based best-model decision.",
            "Explainable assessment: ML risk, lexical risk, highlighted terms, token contribution.",
            "Layered architecture with Streamlit, NLP/ML core, and Supabase/PostgreSQL.",
        ],
        6.9,
        1.75,
        5.55,
        3.8,
        size=17,
    )
    textbox(
        slide,
        "Positioning: decision-support for initial reliability screening, not a replacement for professional fact-checking.",
        1.1,
        6.22,
        11.1,
        0.42,
        size=16,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide)


def add_two_column_slide(prs, number, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, number)
    textbox(slide, left_title, 0.72, 1.22, 5.7, 0.42, size=21, bold=True)
    bullet_list(slide, left_items, 0.9, 1.82, 5.25, 4.6, size=18)
    textbox(slide, right_title, 6.82, 1.22, 5.6, 0.42, size=21, bold=True)
    bullet_list(slide, right_items, 7.0, 1.82, 5.25, 4.6, size=18)
    add_footer(slide)


def add_image_slide(prs, number, title, img_path, caption, left_items=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, number)
    if left_items:
        bullet_list(slide, left_items, 0.75, 1.28, 3.6, 5.45, size=16)
        slide.shapes.add_picture(str(img_path), Inches(4.65), Inches(1.18), width=Inches(7.55))
    else:
        slide.shapes.add_picture(str(img_path), Inches(1.18), Inches(1.16), width=Inches(10.9))
    textbox(slide, caption, 0.75, 6.72, 11.8, 0.34, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(slide)


def add_metric_card(slide, label, value, x, y, w=2.7):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLACK
    textbox(slide, value, x + 0.12, y + 0.13, w - 0.24, 0.35, size=24, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, label, x + 0.12, y + 0.58, w - 0.24, 0.26, size=11, color=GRAY, align=PP_ALIGN.CENTER)


def add_dataset_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Dataset and Label Convention", number)
    textbox(slide, "Primary source", 0.75, 1.18, 4.7, 0.35, size=20, bold=True)
    bullet_list(
        slide,
        [
            "VFND Vietnamese Fake News Dataset",
            "Open research dataset for Vietnamese fake-news studies",
            "Original labels are normalized into reliable vs unreliable",
            "Source link is opened during defense for traceability",
        ],
        0.92,
        1.72,
        5.55,
        2.2,
        size=17,
    )
    textbox(slide, "Source: https://github.com/WhySchools/VFND-vietnamese-fake-news-datasets", 0.92, 4.1, 5.8, 0.4, size=12, color=GRAY)
    add_metric_card(slide, "train samples", "350", 7.0, 1.28)
    add_metric_card(slide, "validation samples", "75", 9.7, 1.28)
    add_metric_card(slide, "test samples", "75", 7.0, 2.6)
    add_metric_card(slide, "label distribution", "251 / 249", 9.7, 2.6)
    textbox(slide, "Label convention", 7.0, 4.25, 5.0, 0.35, size=20, bold=True)
    bullet_list(
        slide,
        [
            "0 = reliable / real news",
            "1 = unreliable / fake / clickbait",
            "Cleaning removed duplicates, missing labels, and very short texts",
        ],
        7.18,
        4.76,
        5.05,
        1.3,
        size=16,
    )
    add_footer(slide)


def add_data_preparation_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Data Preparation Evidence", number)
    add_metric_card(slide, "duplicate ratio removed", "31.32%", 0.8, 1.2, w=2.8)
    add_metric_card(slide, "minimum length threshold", "20", 3.85, 1.2, w=2.8)
    add_metric_card(slide, "full label 0 count", "251", 6.9, 1.2, w=2.8)
    add_metric_card(slide, "full label 1 count", "249", 9.95, 1.2, w=2.8)
    textbox(slide, "Cleaning and validation steps", 0.85, 2.85, 5.7, 0.38, size=21, bold=True)
    bullet_list(
        slide,
        [
            "Remove rows with missing labels or unusable text.",
            "Drop duplicates to reduce leakage and repeated examples.",
            "Keep only records that pass a minimum text-length threshold.",
            "Store processed splits as train.csv, val.csv, and test.csv.",
        ],
        1.05,
        3.42,
        5.4,
        2.1,
        size=17,
    )
    textbox(slide, "Why this matters", 6.95, 2.85, 5.2, 0.38, size=21, bold=True)
    bullet_list(
        slide,
        [
            "A clear split prevents training on the final test data.",
            "Balanced labels make F1 macro and accuracy meaningful.",
            "A reproducible dataset profile supports academic defense.",
            "Limitations are stated honestly because dataset size remains limited.",
        ],
        7.15,
        3.42,
        5.1,
        2.1,
        size=17,
    )
    add_footer(slide)


def add_model_table_slide(prs, number):
    metrics = json.loads((CFG.reports_dir / "metrics_baseline.json").read_text(encoding="utf-8"))
    rows = [
        ("Logistic Regression", metrics["lr"]),
        ("Linear SVM", metrics["svm"]),
        ("Random Forest", metrics["rf"]),
        ("Naive Bayes", metrics["nb"]),
    ]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Model Benchmark", number)
    textbox(slide, "Selection rule: best validation F1 macro, then refit for final test.", 0.78, 1.12, 11.6, 0.35, size=17, color=DARK_GRAY)
    table = slide.shapes.add_table(5, 5, Inches(0.75), Inches(1.75), Inches(11.85), Inches(2.55)).table
    headers = ["Model", "Val Acc", "Val F1 Macro", "Test Acc", "Test F1 Macro"]
    widths = [3.2, 1.55, 2.0, 1.55, 2.1]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                set_run(r, size=12, bold=True, color=WHITE)
    for row_idx, (name, item) in enumerate(rows, start=1):
        values = [
            name,
            f"{item['validation']['accuracy']:.4f}",
            f"{item['validation']['f1_macro']:.4f}",
            f"{item['test']['accuracy']:.4f}",
            f"{item['test']['f1_macro']:.4f}",
        ]
        for col, value in enumerate(values):
            cell = table.cell(row_idx, col)
            cell.text = value
            if name == "Linear SVM":
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if col else PP_ALIGN.LEFT
                for r in p.runs:
                    set_run(r, size=12, bold=(name == "Linear SVM"))
    bullet_list(
        slide,
        [
            "Linear SVM has the strongest validation F1 macro.",
            "Final refit result: Accuracy 0.9200, F1 macro 0.9199, ROC-AUC 0.9915.",
            "Four baselines make the model choice defensible, not arbitrary.",
        ],
        1.0,
        4.72,
        11.1,
        1.2,
        size=18,
    )
    add_footer(slide)


def add_testing_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Testing and Reproducibility", number)
    textbox(slide, "Automated checks", 0.78, 1.15, 5.6, 0.38, size=21, bold=True)
    bullet_list(
        slide,
        [
            "Unit tests: pytest checks preprocessing, inference, and storage fallback behavior.",
            "Artifact evaluation: scripts/evaluate.py verifies required model/report files.",
            "Notebook validation: Colab training notebook is JSON-parseable.",
            "App smoke test: Streamlit returns HTTP 200 on localhost.",
        ],
        0.96,
        1.72,
        5.65,
        3.1,
        size=17,
    )
    textbox(slide, "Defense readiness", 6.9, 1.15, 5.4, 0.38, size=21, bold=True)
    bullet_list(
        slide,
        [
            "Prepared reliable and suspicious cases for stable live demo.",
            "Dataset source and processed profile are opened for traceability.",
            "Metrics and confusion matrix are available before the app demo.",
            "No secret keys are shown or committed in source code.",
        ],
        7.08,
        1.72,
        5.25,
        3.1,
        size=17,
    )
    textbox(slide, "Current verification: 12 tests passed; pipeline evaluation returns ok.", 1.1, 6.15, 11.2, 0.4, size=18, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)


def add_demo_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Defense Demo Flow", number)
    steps = [
        ("1", "Dashboard", "Show dataset evidence, model benchmark, workflow coverage, confusion matrix."),
        ("2", "Reliable Case", "Paste prepared Vietnamese article text; explain low risk and low lexical signals."),
        ("3", "Suspicious Case", "Paste prepared clickbait-like text; point to highlighted terms and lexical risk."),
        ("4", "Feedback + History", "Submit feedback and open history to prove the Supabase workflow."),
    ]
    y = 1.3
    for idx, title, desc in steps:
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.95), Inches(y), Inches(0.55), Inches(0.55))
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLACK
        circ.line.fill.background()
        tf = circ.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = idx
        set_run(r, size=14, bold=True, color=WHITE)
        textbox(slide, title, 1.75, y - 0.02, 2.7, 0.35, size=20, bold=True)
        textbox(slide, desc, 4.35, y + 0.02, 7.8, 0.4, size=16, color=DARK_GRAY)
        y += 1.15
    textbox(slide, "Core message", 0.95, 6.05, 2.2, 0.3, size=18, bold=True)
    textbox(
        slide,
        "The app is not only a text box. It is a complete review workflow: prediction, explanation, storage, feedback, and report export.",
        3.0,
        6.02,
        9.0,
        0.55,
        size=16,
        color=DARK_GRAY,
    )
    add_footer(slide)


def add_references_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "References and Project Artifacts", number)
    textbox(slide, "External references", 0.78, 1.15, 5.8, 0.38, size=21, bold=True)
    bullet_list(
        slide,
        [
            "VFND Vietnamese Fake News Dataset: github.com/WhySchools/VFND-vietnamese-fake-news-datasets",
            "Full Fact AI: fullfact.ai",
            "NewsGuard rating process: newsguardtech.com/ratings/rating-process-criteria",
            "Google Fact Check tools and ClaimReview documentation",
            "ClaimBuster: idir.uta.edu/claimbuster",
        ],
        0.95,
        1.72,
        5.85,
        3.7,
        size=15,
    )
    textbox(slide, "Internal artifacts", 7.0, 1.15, 5.4, 0.38, size=21, bold=True)
    bullet_list(
        slide,
        [
            "reports/dataset_profile.md",
            "reports/model_comparison.md",
            "reports/figures/confusion_matrix_svm.png",
            "notebooks/colab_train_baseline.ipynb",
            "docs/CHUAN_BI_BAO_VE_04_06_VI.md",
        ],
        7.18,
        1.72,
        5.1,
        3.1,
        size=16,
    )
    textbox(
        slide,
        "All metrics and diagrams are generated from the project repository to support reproducibility.",
        1.05,
        6.25,
        11.2,
        0.4,
        size=16,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide)


def build_deck():
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_objectives_slide(prs, 2)
    add_two_column_slide(
        prs,
        3,
        "Problem and Scope",
        "Problem",
        [
            "Vietnamese users face misinformation, clickbait, and low-reliability articles.",
            "A single fake/real label is not enough for user trust.",
            "A practical tool should show risk, evidence-like signals, and limitations.",
        ],
        "Project Scope",
        [
            "Vietnamese text reliability screening.",
            "NLP/ML baseline models with measurable evaluation.",
            "Visual explanation, history, feedback, and case report export.",
        ],
    )
    add_two_column_slide(
        prs,
        4,
        "Benchmark-Inspired Requirements",
        "Observed patterns",
        [
            "Fact-checking tools support reviewers instead of replacing them.",
            "Reliability systems expose ratings plus supporting criteria.",
            "Modern ML tools need traceability, review history, and feedback.",
        ],
        "Implemented in this project",
        [
            "Assessment Summary with risk band and score.",
            "ML risk, lexical risk, suspicious-term highlighting.",
            "Supabase history, user feedback loop, and downloadable case report.",
        ],
    )
    add_dataset_slide(prs, 5)
    add_data_preparation_slide(prs, 6)
    add_image_slide(
        prs,
        7,
        "NLP and ML Pipeline",
        CFG.reports_figures_dir / "report_training_pipeline.png",
        "Dataset cleaning, label normalization, TF-IDF extraction, model training, evaluation, and export.",
    )
    add_two_column_slide(
        prs,
        8,
        "Methodology",
        "Feature extraction",
        [
            "Clean Vietnamese news text.",
            "Represent articles with TF-IDF vectors.",
            "Use sparse text features for efficient baseline learning.",
        ],
        "Models",
        [
            "Logistic Regression: probabilistic linear baseline.",
            "Linear SVM: strong margin-based classifier for TF-IDF.",
            "Random Forest: tree ensemble comparison.",
            "Naive Bayes: fast classical text baseline.",
        ],
    )
    add_model_table_slide(prs, 9)
    add_image_slide(
        prs,
        10,
        "Best Model Evaluation",
        CFG.reports_figures_dir / "confusion_matrix_svm.png",
        "Linear SVM is selected by validation F1 macro; confusion matrix supports transparent error discussion.",
        left_items=[
            "Best model: Linear SVM",
            "Final accuracy: 0.9200",
            "Final F1 macro: 0.9199",
            "Final ROC-AUC: 0.9915",
            "Metric choice: F1 macro handles both labels fairly.",
        ],
    )
    add_image_slide(
        prs,
        11,
        "Layered Architecture",
        CFG.reports_figures_dir / "report_architecture.png",
        "Streamlit UI, NLP/ML Core Engine, and Supabase/PostgreSQL storage are separated for maintainability.",
    )
    add_image_slide(
        prs,
        12,
        "Prediction Workflow",
        CFG.reports_figures_dir / "report_sequence_prediction.png",
        "User input flows through preprocessing, model inference, explanation, persistence, and feedback.",
    )
    add_image_slide(
        prs,
        13,
        "Database Design",
        CFG.reports_figures_dir / "report_database_erd.png",
        "The database stores predictions and feedback so the application can support review history and future retraining.",
        left_items=[
            "Prediction records keep model output and risk scores.",
            "Feedback records user judgment and notes.",
            "History supports reviewer inspection.",
            "Schema separates app workflow from ML artifacts.",
        ],
    )
    add_image_slide(
        prs,
        14,
        "Database and Feedback Loop",
        CFG.reports_figures_dir / "report_feedback_loop.png",
        "Prediction history and user feedback support reviewer workflow and future retraining.",
        left_items=[
            "Predictions are stored for audit/history.",
            "Feedback records user judgment.",
            "Reviewer can inspect mistakes.",
            "Future training can reuse validated feedback.",
        ],
    )
    add_image_slide(
        prs,
        15,
        "Result Dashboard and Explainability",
        CFG.reports_figures_dir / "report_result_dashboard.png",
        "The result page presents a readable risk summary, chart, highlighted input, and token-level explanation.",
    )
    add_demo_slide(prs, 16)
    add_testing_slide(prs, 17)
    add_two_column_slide(
        prs,
        18,
        "Limitations and Future Work",
        "Current limitations",
        [
            "Dataset size is still limited.",
            "The tool screens linguistic reliability signals, not absolute truth.",
            "No claim-level evidence retrieval yet.",
            "Source credibility is not fully modeled.",
        ],
        "Future work",
        [
            "Expand and continuously update Vietnamese news datasets.",
            "Fine-tune PhoBERT or other Vietnamese language models.",
            "Add source credibility and claim evidence retrieval.",
            "Use feedback data for monitored retraining.",
        ],
    )
    add_two_column_slide(
        prs,
        19,
        "Conclusion",
        "What was built",
        [
            "A working Streamlit web app for Vietnamese news reliability screening.",
            "A reproducible NLP/ML pipeline with four baseline models.",
            "Visual explanation, history, feedback, and report export.",
        ],
        "Defense message",
        [
            "The project is a decision-support tool, not an absolute fact-checker.",
            "The model choice is justified by metrics, not preference.",
            "The architecture can be extended to stronger models and richer evidence.",
        ],
    )
    add_references_slide(prs, 20)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build_deck()
